import streamlit as st
import os
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pdf2image import convert_from_path
import easyocr
import tempfile
from PIL import Image

# ==========================================
# 1. Caching AI Model (โหลดครั้งเดียว)
# ==========================================
@st.cache_resource(show_spinner=False)
def load_easyocr_reader(langs):
    return easyocr.Reader(list(langs), gpu=False)

# ==========================================
# 2. Slide Converter Class
# ==========================================
class SlideConverterAI:
    def __init__(self, reader):
        self.reader = reader
        self.temp_files = []

    def cleanup_temp_files(self):
        for f in self.temp_files:
            try:
                if os.path.exists(f): os.remove(f)
            except: pass
        self.temp_files.clear()

    def get_smart_fill_color(self, img, bbox):
        x, y, w, h = bbox
        h_img, w_img, _ = img.shape
        margin = 5
        pixels = []

        if y > margin: pixels.extend(img[y-margin:y, x:x+w].reshape(-1, 3))
        if y+h+margin < h_img: pixels.extend(img[y+h:y+h+margin, x:x+w].reshape(-1, 3))
        if x > margin: pixels.extend(img[y:y+h, x-margin:x].reshape(-1, 3))
        if x+w+margin < w_img: pixels.extend(img[y:y+h, x+w:x+w+margin].reshape(-1, 3))

        if pixels:
            bgr = np.median(np.array(pixels), axis=0)
            return int(bgr[2]), int(bgr[1]), int(bgr[0])
        return 255, 255, 255

    def remove_text_from_image(self, img, text_bboxes):
        mask = np.zeros(img.shape[:2], np.uint8)
        for x, y, w, h in text_bboxes:
            cv2.rectangle(mask, (x-3, y-3), (x+w+3, y+h+3), 255, -1)
        return cv2.inpaint(img, mask, 3, cv2.INPAINT_TELEA) if np.any(mask) else img

    def convert_to_editable_pptx(self, input_path, output_path, progress_callback=None):
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        try:
            if input_path.lower().endswith(".pdf"):
                images = convert_from_path(input_path) 
            else:
                images = [Image.open(input_path)]

            total_slides = len(images)
            for idx, pil_img in enumerate(images):
                if progress_callback:
                    # คำนวณ % ความคืบหน้า (ให้เหลือ 10% ไว้ตอนเซฟไฟล์)
                    progress = int((idx / total_slides) * 90)
                    progress_callback(f"กำลังประมวลผลสไลด์ {idx+1}/{total_slides}...", progress)

                cv_img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
                h, w, _ = cv_img.shape
                sx = prs.slide_width.inches / w
                sy = prs.slide_height.inches / h

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                ocr = self.reader.readtext(cv_img)
                boxes = []

                for (bbox, text, conf) in ocr:
                    if conf < 0.3: continue
                    tl, _, br, _ = bbox
                    x, y = int(tl[0]), int(tl[1])
                    bw, bh = int(br[0]-tl[0]), int(br[1]-tl[1])
                    if bw > 10 and bh > 5:
                        boxes.append((x, y, bw, bh, text))

                bg = self.remove_text_from_image(cv_img, [(x,y,bw,bh) for x,y,bw,bh,_ in boxes])
                
                # ใช้ tempfile สำหรับรูปภาพแต่ละสไลด์
                tmp_fd, tmp_path = tempfile.mkstemp(suffix=".jpg")
                os.close(tmp_fd)
                self.temp_files.append(tmp_path)
                cv2.imwrite(tmp_path, bg)

                slide.shapes.add_picture(tmp_path, 0, 0, prs.slide_width, prs.slide_height)

                for x,y,bw,bh,text in boxes:
                    box = slide.shapes.add_textbox(Inches(x*sx), Inches(y*sy), Inches(bw*sx), Inches(bh*sy))
                    tf = box.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = text
                    p.alignment = PP_ALIGN.LEFT

                    size = max(10, min(40, bh*sy*10))
                    p.font.size = Pt(int(size))
                    p.font.name = "Arial"

                    r,g,b = self.get_smart_fill_color(cv_img, (x,y,bw,bh))
                    brightness = r*0.299 + g*0.587 + b*0.114
                    p.font.color.rgb = RGBColor(0,0,0) if brightness>128 else RGBColor(255,255,255)

            if progress_callback:
                progress_callback("กำลังสร้างไฟล์ PowerPoint...", 95)
            
            prs.save(output_path)
            
            if progress_callback:
                progress_callback("✅ แปลงไฟล์สำเร็จ!", 100)

        finally:
            self.cleanup_temp_files()


# ==========================================
# 3. Streamlit UI
# ==========================================
st.set_page_config(page_title="PDF to Editable PPTX", page_icon="📽️", layout="centered")

st.title("📽️ Slide to Editable PPTX")
st.markdown("แปลงไฟล์ PDF หรือรูปภาพสไลด์ ให้เป็น **PowerPoint ที่สามารถแก้ไขข้อความได้** ด้วย AI (EasyOCR)")

# เลือกภาษา
lang_option = st.selectbox("🌐 เลือกภาษาในสไลด์", ["ไทย + อังกฤษ", "ไทยเท่านั้น", "อังกฤษเท่านั้น"])
lang_map = {
    "ไทย + อังกฤษ": ("th", "en"),
    "ไทยเท่านั้น": ("th",),
    "อังกฤษเท่านั้น": ("en",)
}

# โหลดโมเดล
with st.spinner("⏳ กำลังโหลดโมเดล AI... (ครั้งแรกอาจใช้เวลาสักครู่)"):
    reader = load_easyocr_reader(lang_map[lang_option])

# อัปโหลดไฟล์
uploaded_file = st.file_uploader("📂 อัปโหลดไฟล์ PDF หรือรูปภาพ", type=['pdf', 'png', 'jpg', 'jpeg'])

if uploaded_file is not None:
    st.info(f"📄 ไฟล์ที่เลือก: {uploaded_file.name}")
    
    if st.button("✨ เริ่มทำการแปลงไฟล์", type="primary"):
        # สร้าง Layout สำหรับแสดงผลความคืบหน้า
        progress_bar = st.progress(0)
        status_text = st.empty()

        def st_progress_callback(msg, val):
            status_text.text(msg)
            progress_bar.progress(val)

        # จัดการไฟล์ Temp สำหรับ Input และ Output
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{uploaded_file.name.split('.')[-1]}") as tmp_in:
            tmp_in.write(uploaded_file.getvalue())
            input_path = tmp_in.name
            
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_out:
            output_path = tmp_out.name

        try:
            # เริ่มทำงาน
            converter = SlideConverterAI(reader)
            converter.convert_to_editable_pptx(input_path, output_path, progress_callback=st_progress_callback)
            
            # เมื่อเสร็จแล้ว ให้ดาวน์โหลด
            with open(output_path, "rb") as f:
                pptx_data = f.read()
                
            st.success("🎉 แปลงไฟล์เสร็จสมบูรณ์!")
            
            st.download_button(
                label="📥 ดาวน์โหลดไฟล์ PowerPoint (.pptx)",
                data=pptx_data,
                file_name=f"{os.path.splitext(uploaded_file.name)[0]}_editable.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )
            
        except Exception as e:
            st.error(f"❌ เกิดข้อผิดพลาด: {str(e)}")
            
        finally:
            # ลบไฟล์ Temp ทิ้งเมื่อเสร็จสิ้น
            if os.path.exists(input_path): os.remove(input_path)
            if os.path.exists(output_path): os.remove(output_path)