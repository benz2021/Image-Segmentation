# main.py (Auto Install & Run Version)
import sys
import subprocess
import importlib
import os
import time

# ==============================
# AUTO INSTALL LIBRARIES
# ==============================
REQUIRED_PACKAGES = {
    "cv2": "opencv-python",
    "PIL": "Pillow",
    "pptx": "python-pptx",
    "numpy": "numpy<2.0.0",
    "pdf2image": "pdf2image",
    "easyocr": "easyocr",
    "torch": "torch"
}

def install_package(pkg):
    print(f"📦 Installing {pkg} ...")
    subprocess.check_call(
        [sys.executable, "-m", "pip", "install", pkg],
        stdout=subprocess.DEVNULL
    )

def ensure_libraries():
    print("🔍 Checking libraries...")
    for module, package in REQUIRED_PACKAGES.items():
        try:
            importlib.import_module(module)
            print(f"✅ {package}")
        except ImportError:
            install_package(package)

    print("🎉 All libraries ready!\n")

# ==============================
# POPPLER CHECK (WINDOWS)
# ==============================
def check_poppler():
    if sys.platform != "win32":
        return

    try:
        from pdf2image import convert_from_path
        convert_from_path  # test import only
    except Exception:
        print("""
❌ ไม่พบ Poppler (จำเป็นสำหรับ PDF)

วิธีแก้ (เลือกอย่างใดอย่างหนึ่ง):
1) ดาวน์โหลด:
   https://github.com/oschwartz10612/poppler-windows/releases/
   แล้วเพิ่ม path:  ...\\poppler\\Library\\bin

หรือ
2) ถ้ามี Chocolatey:
   choco install poppler

แล้วรันโปรแกรมใหม่อีกครั้ง
""")
        input("กด Enter เพื่อออก...")
        sys.exit(1)

# ==============================
# MAIN PROGRAM
# ==============================
def main():
    print("""
╔══════════════════════════════════════════════════════════╗
║     NotebookLM Slides to Editable PPTX Converter v2.0    ║
║              (Codia.ai Style Clone)                      ║
╚══════════════════════════════════════════════════════════╝
""")

    # 1. ตรวจ + ติดตั้ง lib อัตโนมัติ
    ensure_libraries()

    # 2. ตรวจ poppler
    check_poppler()

    # 3. เรียก UI
    try:
        from ui_main import CodiaCloneApp
        print("🚀 Starting application...\n")
        app = CodiaCloneApp()
        app.run()

    except Exception as e:
        print("❌ Program Error:", e)
        import traceback
        traceback.print_exc()
        input("\nกด Enter เพื่อออก...")

if __name__ == "__main__":
    main()

# ui_main.py (แก้ไขแล้ว)
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
import threading
import os
import sys
import subprocess
import easyocr  # เพิ่ม import

# Import ตัวแปลงที่เราสร้าง
from slide_converter import SlideConverterAI

class CodiaCloneApp:
    def __init__(self):
        self.root = tk.Tk()
        self.root.title("NotebookLM Slides to Editable PPTX (Codia Style)")
        self.root.geometry("800x550")
        self.root.configure(bg="#f4f6f9")
        self.root.resizable(False, False)
        
        # Set icon if available
        try:
            self.root.iconbitmap(default='icon.ico')
        except:
            pass
            
        self.converter = SlideConverterAI()
        self.input_file = None
        self.conversion_running = False
        
        self.setup_ui()
        
        # Load Model ใน background
        self.status_lbl.config(text="⏳ Initializing AI Models...")
        threading.Thread(target=self.init_model, daemon=True).start()

    def init_model(self):
        """โหลดโมเดล AI"""
        try:
            self.converter.load_model()
            self.root.after(0, self.on_model_loaded)
        except Exception as e:
            self.root.after(0, lambda: self.on_model_error(str(e)))

    def on_model_loaded(self):
        """เมื่อโหลดโมเดลสำเร็จ"""
        self.status_lbl.config(text="✅ Ready to Convert")
        self.convert_btn.config(state=tk.NORMAL)
        self.browse_btn.config(state=tk.NORMAL)

    def on_model_error(self, error_msg):
        """เมื่อโหลดโมเดลไม่สำเร็จ"""
        self.status_lbl.config(text=f"❌ Model Error: {error_msg[:50]}...")
        messagebox.showerror("Model Error", 
                            f"ไม่สามารถโหลด AI Model ได้:\n{error_msg}\n\n"
                            "กรุณาตรวจสอบการติดตั้ง:\n"
                            "pip install easyocr torch")

    def setup_ui(self):
        """สร้าง UI"""
        # Header
        header = tk.Frame(self.root, bg="white", pady=20)
        header.pack(fill=tk.X)
        
        tk.Label(header, text="📽️ NotebookLM Slides to PowerPoint", 
                font=("Helvetica", 18, "bold"), bg="white", fg="#2c3e50").pack()
        tk.Label(header, text="แปลง PDF/ภาพเป็น PowerPoint ที่แก้ไขได้ด้วย AI", 
                font=("Helvetica", 10), bg="white", fg="#7f8c8d").pack()

        # Main Content
        main_frame = tk.Frame(self.root, bg="#f4f6f9", pady=20)
        main_frame.pack(fill=tk.BOTH, expand=True)

        # File Selection Area
        self.file_frame = tk.Frame(main_frame, bg="white", bd=2, relief=tk.GROOVE, height=150)
        self.file_frame.pack(pady=10, padx=50, ipadx=20, ipady=20, fill=tk.X)
        self.file_frame.pack_propagate(False)
        
        self.lbl_instruction = tk.Label(self.file_frame, 
                                       text="👇 เลือกไฟล์ PDF หรือภาพที่ต้องการแปลง\n(PDF, JPG, PNG)", 
                                       bg="white", font=("Arial", 11), fg="#34495e")
        self.lbl_instruction.pack(pady=10)
        
        # ปุ่ม Browse
        btn_frame = tk.Frame(self.file_frame, bg="white")
        btn_frame.pack(pady=10)
        
        self.browse_btn = tk.Button(btn_frame, text="📂 เลือกไฟล์", command=self.browse_file,
                                   bg="#3498db", fg="white", font=("Arial", 11, "bold"), 
                                   padx=20, pady=8, bd=0, cursor="hand2", state=tk.DISABLED)
        self.browse_btn.pack(side=tk.LEFT, padx=5)
        
        # ปุ่ม Clear
        self.clear_btn = tk.Button(btn_frame, text="🗑️ ล้าง", command=self.clear_selection,
                                  bg="#95a5a6", fg="white", font=("Arial", 11, "bold"),
                                  padx=20, pady=8, bd=0, cursor="hand2", state=tk.DISABLED)
        self.clear_btn.pack(side=tk.LEFT, padx=5)

        # Options Frame
        opts_frame = tk.Frame(main_frame, bg="#f4f6f9")
        opts_frame.pack(pady=15)
        
        # ตัวเลือกภาษา
        lang_frame = tk.Frame(opts_frame, bg="#f4f6f9")
        lang_frame.pack(pady=5)
        
        tk.Label(lang_frame, text="🌐 ภาษา:", bg="#f4f6f9", font=("Arial", 10)).pack(side=tk.LEFT, padx=5)
        
        self.lang_var = tk.StringVar(value="ไทย+อังกฤษ")
        lang_options = ttk.Combobox(lang_frame, textvariable=self.lang_var, 
                                    values=["ไทย+อังกฤษ", "ไทย", "อังกฤษ"],
                                    state="readonly", width=15)
        lang_options.pack(side=tk.LEFT)
        lang_options.bind('<<ComboboxSelected>>', self.on_language_change)

        # Progress Bar
        self.progress = ttk.Progressbar(main_frame, orient=tk.HORIZONTAL, length=500, mode='determinate')
        self.progress.pack(pady=10)

        # Status Label
        self.status_lbl = tk.Label(main_frame, text="⏳ กำลังโหลด AI Model...", 
                                   bg="#f4f6f9", fg="#7f8c8d", font=("Arial", 10))
        self.status_lbl.pack()

        # Convert Button
        self.convert_btn = tk.Button(main_frame, text="✨ แปลงไฟล์ ✨", command=self.start_conversion,
                                     bg="#27ae60", fg="white", font=("Arial", 14, "bold"), 
                                     padx=40, pady=15, bd=0, cursor="hand2", state=tk.DISABLED)
        self.convert_btn.pack(pady=20)

        # Footer
        footer = tk.Frame(self.root, bg="#ecf0f1", height=30)
        footer.pack(fill=tk.X, side=tk.BOTTOM)
        
        tk.Label(footer, text="Powered by EasyOCR | รองรับภาษาไทย", 
                bg="#ecf0f1", fg="#7f8c8d", font=("Arial", 9)).pack(pady=5)

    def on_language_change(self, event=None):
        """เมื่อเปลี่ยนภาษา"""
        lang_map = {
            "ไทย+อังกฤษ": ['th', 'en'],
            "ไทย": ['th'],
            "อังกฤษ": ['en']
        }
        selected = self.lang_var.get()
        if selected in lang_map:
            # ถ้าต้องการเปลี่ยนภาษาจริงๆ ต้องโหลดโมเดลใหม่
            if messagebox.askyesno("เปลี่ยนภาษา", 
                                   "ต้องโหลดโมเดลใหม่เพื่อเปลี่ยนภาษา\nดำเนินการต่อ?"):
                self.status_lbl.config(text="⏳ กำลังโหลดโมเดลภาษาใหม่...")
                self.convert_btn.config(state=tk.DISABLED)
                threading.Thread(target=self.reload_model, args=(lang_map[selected],), daemon=True).start()

    def reload_model(self, langs):
        """โหลดโมเดลใหม่พร้อมภาษา"""
        try:
            self.converter.reader = easyocr.Reader(langs, gpu=False)
            self.root.after(0, self.on_model_loaded)
        except Exception as e:
            self.root.after(0, lambda: self.on_model_error(str(e)))

    def browse_file(self):
        """เลือกไฟล์"""
        f = filedialog.askopenfilename(
            title="เลือกไฟล์ PDF หรือภาพ",
            filetypes=[
                ("ไฟล์ที่รองรับ", "*.pdf *.png *.jpg *.jpeg *.bmp *.tiff"),
                ("PDF files", "*.pdf"),
                ("Images", "*.png *.jpg *.jpeg *.bmp *.tiff"),
                ("All files", "*.*")
            ]
        )
        if f:
            self.input_file = f
            name = os.path.basename(f)
            self.lbl_instruction.config(text=f"📄 {name}", fg="#27ae60", font=("Arial", 10, "bold"))
            self.clear_btn.config(state=tk.NORMAL)

    def clear_selection(self):
        """ล้างการเลือกไฟล์"""
        self.input_file = None
        self.lbl_instruction.config(text="👇 เลือกไฟล์ PDF หรือภาพที่ต้องการแปลง\n(PDF, JPG, PNG)", 
                                   fg="#34495e", font=("Arial", 11))
        self.clear_btn.config(state=tk.DISABLED)

    def start_conversion(self):
        """เริ่มการแปลง"""
        if not self.input_file:
            messagebox.showwarning("Warning", "กรุณาเลือกไฟล์ก่อน")
            return
        
        if self.conversion_running:
            messagebox.showinfo("Info", "กำลังแปลงไฟล์อยู่ กรุณารอสักครู่")
            return
        
        # เลือกที่บันทึก
        default_name = os.path.splitext(os.path.basename(self.input_file))[0] + "_editable.pptx"
        save_path = filedialog.asksaveasfilename(
            title="บันทึกไฟล์ PowerPoint",
            defaultextension=".pptx",
            initialfile=default_name,
            filetypes=[("PowerPoint", "*.pptx"), ("All files", "*.*")]
        )
        
        if not save_path:
            return

        # ปิดปุ่มต่างๆ
        self.conversion_running = True
        self.convert_btn.config(state=tk.DISABLED)
        self.browse_btn.config(state=tk.DISABLED)
        self.clear_btn.config(state=tk.DISABLED)
        self.progress['value'] = 0
        
        # รันใน thread แยก
        thread = threading.Thread(target=self.run_conversion, 
                                 args=(self.input_file, save_path),
                                 daemon=True)
        thread.start()

    def run_conversion(self, input_path, output_path):
        """รันการแปลงใน thread แยก"""
        try:
            # เรียกใช้ converter
            self.converter.convert_to_editable_pptx(
                input_path, 
                output_path,
                progress_callback=self.update_progress
            )
            
            # แจ้งผลสำเร็จ
            self.root.after(0, lambda: self.on_conversion_success(output_path))
            
        except Exception as e:
            self.root.after(0, lambda: self.on_conversion_error(str(e)))
        
        finally:
            self.conversion_running = False

    def update_progress(self, message, value):
        """อัปเดตความคืบหน้า"""
        self.root.after(0, lambda: self._update_ui(message, value))

    def _update_ui(self, message, value):
        """อัปเดต UI"""
        self.status_lbl.config(text=message)
        self.progress['value'] = value

    def on_conversion_success(self, output_path):
        """เมื่อแปลงสำเร็จ"""
        self.status_lbl.config(text="✅ แปลงไฟล์สำเร็จ!")
        self.progress['value'] = 100
        
        # ถามว่าเปิดไฟล์หรือไม่
        if messagebox.askyesno("Success", f"บันทึกไฟล์ที่:\n{output_path}\n\nต้องการเปิดไฟล์หรือไม่?"):
            self.open_file(output_path)
        
        # เปิดปุ่มอีกครั้ง
        self.convert_btn.config(state=tk.NORMAL)
        self.browse_btn.config(state=tk.NORMAL)

    def on_conversion_error(self, error_msg):
        """เมื่อเกิดข้อผิดพลาด"""
        self.status_lbl.config(text=f"❌ Error: {error_msg[:50]}...")
        self.progress['value'] = 0
        
        messagebox.showerror("Error", f"เกิดข้อผิดพลาด:\n{error_msg}")
        
        # เปิดปุ่มอีกครั้ง
        self.convert_btn.config(state=tk.NORMAL)
        self.browse_btn.config(state=tk.NORMAL)
        self.clear_btn.config(state=tk.NORMAL if self.input_file else tk.DISABLED)

    def open_file(self, filepath):
        """เปิดไฟล์ด้วยโปรแกรมเริ่มต้น"""
        try:
            if sys.platform == 'win32':
                os.startfile(filepath)
            elif sys.platform == 'darwin':
                subprocess.run(['open', filepath])
            else:
                subprocess.run(['xdg-open', filepath])
        except Exception as e:
            print(f"ไม่สามารถเปิดไฟล์ได้: {e}")

    def run(self):
        """เริ่มโปรแกรม"""
        self.root.mainloop()

if __name__ == "__main__":
    app = CodiaCloneApp()
    app.run()

import os
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pdf2image import convert_from_path
import easyocr
import tempfile
from PIL import Image

class SlideConverterAI:
    def __init__(self):
        self.reader = None
        self.is_loaded = False
        self.temp_files = []

    # ---------- MODEL ----------
    def load_model(self, langs=('th', 'en')):
        if self.is_loaded:
            return
        print("Loading EasyOCR model...")
        self.reader = easyocr.Reader(list(langs), gpu=False)
        self.is_loaded = True
        print("✅ Model loaded")

    # ---------- CLEANUP ----------
    def cleanup_temp_files(self):
        for f in self.temp_files:
            try:
                if os.path.exists(f):
                    os.remove(f)
            except:
                pass
        self.temp_files.clear()

    # ---------- COLOR ----------
    def get_smart_fill_color(self, img, bbox):
        x, y, w, h = bbox
        h_img, w_img, _ = img.shape
        margin = 5
        pixels = []

        if y > margin:
            pixels.extend(img[y-margin:y, x:x+w].reshape(-1, 3))
        if y+h+margin < h_img:
            pixels.extend(img[y+h:y+h+margin, x:x+w].reshape(-1, 3))
        if x > margin:
            pixels.extend(img[y:y+h, x-margin:x].reshape(-1, 3))
        if x+w+margin < w_img:
            pixels.extend(img[y:y+h, x+w:x+w+margin].reshape(-1, 3))

        if pixels:
            bgr = np.median(np.array(pixels), axis=0)
            return int(bgr[2]), int(bgr[1]), int(bgr[0])
        return 255, 255, 255

    # ---------- IMAGE ----------
    def remove_text_from_image(self, img, text_bboxes):
        mask = np.zeros(img.shape[:2], np.uint8)
        for x, y, w, h in text_bboxes:
            cv2.rectangle(mask, (x-3, y-3), (x+w+3, y+h+3), 255, -1)
        return cv2.inpaint(img, mask, 3, cv2.INPAINT_TELEA) if np.any(mask) else img

    # ---------- MAIN ----------
    def convert_to_editable_pptx(self, input_path, output_path, progress_callback=None):
        if not self.is_loaded:
            self.load_model()

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        poppler = os.getenv("POPPLER_PATH")

        try:
            if input_path.lower().endswith(".pdf"):
                images = convert_from_path(input_path, poppler_path=poppler)
            else:
                images = [Image.open(input_path)]

            for idx, pil_img in enumerate(images):
                if progress_callback:
                    progress_callback(f"Processing slide {idx+1}/{len(images)}", idx/len(images)*80)

                cv_img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
                h, w, _ = cv_img.shape
                sx = prs.slide_width.inches / w
                sy = prs.slide_height.inches / h

                slide = prs.slides.add_slide(prs.slide_layouts[6])

                ocr = self.reader.readtext(cv_img)
                boxes = []

                for (bbox, text, conf) in ocr:
                    if conf < 0.3:
                        continue
                    tl, _, br, _ = bbox
                    x, y = int(tl[0]), int(tl[1])
                    bw, bh = int(br[0]-tl[0]), int(br[1]-tl[1])
                    if bw > 10 and bh > 5:
                        boxes.append((x, y, bw, bh, text))

                bg = self.remove_text_from_image(cv_img, [(x,y,bw,bh) for x,y,bw,bh,_ in boxes])
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg").name
                self.temp_files.append(tmp)
                cv2.imwrite(tmp, bg)

                slide.shapes.add_picture(tmp, 0, 0, prs.slide_width, prs.slide_height)

                for x,y,bw,bh,text in boxes:
                    box = slide.shapes.add_textbox(
                        Inches(x*sx), Inches(y*sy),
                        Inches(bw*sx), Inches(bh*sy)
                    )
                    tf = box.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = text
                    p.alignment = PP_ALIGN.LEFT

                    size = max(10, min(40, bh*sy*10))
                    p.font.size = Pt(int(size))
                    p.font.name = "Arial"

                    r,g,b = self.get_smart_fill_color(cv_img, (x,y,bw,bh))
                    brightness = r*0.299 + g*0.587 + b*0.114
                    p.font.color.rgb = RGBColor(0,0,0) if brightness>128 else RGBColor(255,255,255)

            prs.save(output_path)
            if progress_callback:
                progress_callback("✅ Done", 100)

        finally:
            self.cleanup_temp_files()
