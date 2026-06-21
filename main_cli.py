import os
import cv2
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pdf2image import convert_from_path
import easyocr
import tempfile
from PIL import Image
import argparse

class SlideConverterAI:
    # (ใช้โค้ด SlideConverterAI เดิมของคุณทั้งหมดที่นี่)
    def __init__(self):
        self.reader = None
        self.is_loaded = False
        self.temp_files = []

    def load_model(self, langs=('th', 'en')):
        if self.is_loaded: return
        print("⏳ Loading EasyOCR model...")
        self.reader = easyocr.Reader(list(langs), gpu=False)
        self.is_loaded = True
        print("✅ Model loaded")

    def cleanup_temp_files(self):
        for f in self.temp_files:
            try:
                if os.path.exists(f): os.remove(f)
            except: pass
        self.temp_files.clear()

    def get_smart_fill_color(self, img, bbox):
        x, y, w, h = bbox
        h_img, w_img, _ = img.shape
        margin = 5
        pixels = []

        if y > margin: pixels.extend(img[y-margin:y, x:x+w].reshape(-1, 3))
        if y+h+margin < h_img: pixels.extend(img[y+h:y+h+margin, x:x+w].reshape(-1, 3))
        if x > margin: pixels.extend(img[y:y+h, x-margin:x].reshape(-1, 3))
        if x+w+margin < w_img: pixels.extend(img[y:y+h, x+w:x+w+margin].reshape(-1, 3))

        if pixels:
            bgr = np.median(np.array(pixels), axis=0)
            return int(bgr[2]), int(bgr[1]), int(bgr[0])
        return 255, 255, 255

    def remove_text_from_image(self, img, text_bboxes):
        mask = np.zeros(img.shape[:2], np.uint8)
        for x, y, w, h in text_bboxes:
            cv2.rectangle(mask, (x-3, y-3), (x+w+3, y+h+3), 255, -1)
        return cv2.inpaint(img, mask, 3, cv2.INPAINT_TELEA) if np.any(mask) else img

    def convert_to_editable_pptx(self, input_path, output_path):
        if not self.is_loaded: self.load_model()

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        try:
            print(f"📄 Processing: {input_path}")
            if input_path.lower().endswith(".pdf"):
                images = convert_from_path(input_path) # ไม่ต้องพึ่ง POPPLER_PATH บน Ubuntu
            else:
                images = [Image.open(input_path)]

            for idx, pil_img in enumerate(images):
                print(f"🔄 Converting slide {idx+1}/{len(images)}...")
                cv_img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)
                h, w, _ = cv_img.shape
                sx = prs.slide_width.inches / w
                sy = prs.slide_height.inches / h

                slide = prs.slides.add_slide(prs.slide_layouts[6])
                ocr = self.reader.readtext(cv_img)
                boxes = []

                for (bbox, text, conf) in ocr:
                    if conf < 0.3: continue
                    tl, _, br, _ = bbox
                    x, y = int(tl[0]), int(tl[1])
                    bw, bh = int(br[0]-tl[0]), int(br[1]-tl[1])
                    if bw > 10 and bh > 5:
                        boxes.append((x, y, bw, bh, text))

                bg = self.remove_text_from_image(cv_img, [(x,y,bw,bh) for x,y,bw,bh,_ in boxes])
                tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".jpg").name
                self.temp_files.append(tmp)
                cv2.imwrite(tmp, bg)

                slide.shapes.add_picture(tmp, 0, 0, prs.slide_width, prs.slide_height)

                for x,y,bw,bh,text in boxes:
                    box = slide.shapes.add_textbox(Inches(x*sx), Inches(y*sy), Inches(bw*sx), Inches(bh*sy))
                    tf = box.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = text
                    p.alignment = PP_ALIGN.LEFT

                    size = max(10, min(40, bh*sy*10))
                    p.font.size = Pt(int(size))
                    p.font.name = "Arial"

                    r,g,b = self.get_smart_fill_color(cv_img, (x,y,bw,bh))
                    brightness = r*0.299 + g*0.587 + b*0.114
                    p.font.color.rgb = RGBColor(0,0,0) if brightness>128 else RGBColor(255,255,255)

            prs.save(output_path)
            print(f"🎉 Success! Saved to {output_path}")

        finally:
            self.cleanup_temp_files()

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert PDF/Image to Editable PPTX via AI")
    parser.add_argument("--input", required=True, help="Path to input PDF or Image")
    parser.add_argument("--output", required=True, help="Path to output PPTX file")
    args = parser.parse_args()

    converter = SlideConverterAI()
    converter.convert_to_editable_pptx(args.input, args.output)